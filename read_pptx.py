import sys
import subprocess

try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
    from pptx import Presentation

def extract_text(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
        print("\n")

if __name__ == "__main__":
    extract_text(r"c:\Workstation\game\게임 장르의 이해.pptx")
